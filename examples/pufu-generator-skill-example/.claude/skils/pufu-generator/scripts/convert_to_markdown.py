#!/usr/bin/env python3
"""
ドキュメントをマークダウン形式に変換するスクリプト

対応形式: .pptx, .xlsx, .pdf, .docx, .txt, .md
出力: マークダウンファイル
"""

import os
import sys
import json
from pathlib import Path
from datetime import datetime


def convert_pptx_to_markdown(filepath: str) -> str:
    """PowerPointをマークダウンに変換"""
    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx をインストールしてください: pip install python-pptx --break-system-packages")
        sys.exit(1)
    
    prs = Presentation(filepath)
    md_lines = [f"# {Path(filepath).stem}\n"]
    md_lines.append(f"*変換元: {Path(filepath).name}*\n")
    
    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"## スライド {i}")
        
        if slide.slide_layout and slide.slide_layout.name:
            md_lines.append(f"*レイアウト: {slide.slide_layout.name}*\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if shape.shape_type == 14:  # TITLE
                    md_lines.append(f"**{text}**\n")
                else:
                    md_lines.append(f"{text}\n")
        
        md_lines.append("")
    
    return "\n".join(md_lines)


def convert_xlsx_to_markdown(filepath: str) -> str:
    """Excelをマークダウンに変換"""
    try:
        import openpyxl
    except ImportError:
        print("openpyxl をインストールしてください: pip install openpyxl --break-system-packages")
        sys.exit(1)
    
    wb = openpyxl.load_workbook(filepath, data_only=True)
    md_lines = [f"# {Path(filepath).stem}\n"]
    md_lines.append(f"*変換元: {Path(filepath).name}*\n")
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        md_lines.append(f"## シート: {sheet_name}\n")
        
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            md_lines.append("*データなし*\n")
            continue
        
        non_empty_rows = [row for row in rows if any(cell is not None for cell in row)]
        if not non_empty_rows:
            md_lines.append("*データなし*\n")
            continue
        
        max_cols = max(len(row) for row in non_empty_rows)
        
        header = non_empty_rows[0]
        header_cells = [str(cell) if cell is not None else "" for cell in header[:max_cols]]
        md_lines.append("| " + " | ".join(header_cells) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(header_cells)) + " |")
        
        for row in non_empty_rows[1:]:
            cells = [str(cell) if cell is not None else "" for cell in row[:max_cols]]
            while len(cells) < len(header_cells):
                cells.append("")
            md_lines.append("| " + " | ".join(cells) + " |")
        
        md_lines.append("")
    
    return "\n".join(md_lines)


def convert_pdf_to_markdown(filepath: str) -> str:
    """PDFをマークダウンに変換"""
    md_lines = [f"# {Path(filepath).stem}\n"]
    md_lines.append(f"*変換元: {Path(filepath).name}*\n")
    
    # pdfplumberを優先、なければPyPDF2を使用
    try:
        import pdfplumber
        
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                md_lines.append(f"## ページ {i}\n")
                
                # テキスト抽出
                text = page.extract_text()
                if text:
                    md_lines.append(text)
                
                # テーブル抽出を試みる
                tables = page.extract_tables()
                for table in tables:
                    if table and len(table) > 0:
                        md_lines.append("\n### テーブル\n")
                        # ヘッダー
                        header = [str(cell) if cell else "" for cell in table[0]]
                        md_lines.append("| " + " | ".join(header) + " |")
                        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                        # データ行
                        for row in table[1:]:
                            cells = [str(cell) if cell else "" for cell in row]
                            md_lines.append("| " + " | ".join(cells) + " |")
                
                md_lines.append("")
        
        return "\n".join(md_lines)
    
    except ImportError:
        pass
    
    # PyPDF2をフォールバックとして使用
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(filepath)
        for i, page in enumerate(reader.pages, 1):
            md_lines.append(f"## ページ {i}\n")
            text = page.extract_text()
            if text:
                md_lines.append(text)
            md_lines.append("")
        
        return "\n".join(md_lines)
    
    except ImportError:
        print("PDF処理ライブラリをインストールしてください:")
        print("  pip install pdfplumber --break-system-packages")
        print("  または")
        print("  pip install PyPDF2 --break-system-packages")
        sys.exit(1)


def convert_docx_to_markdown(filepath: str) -> str:
    """Wordをマークダウンに変換"""
    try:
        from docx import Document
    except ImportError:
        print("python-docx をインストールしてください: pip install python-docx --break-system-packages")
        sys.exit(1)
    
    doc = Document(filepath)
    md_lines = [f"# {Path(filepath).stem}\n"]
    md_lines.append(f"*変換元: {Path(filepath).name}*\n")
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            md_lines.append("")
            continue
        
        style_name = para.style.name if para.style else ""
        if "Heading 1" in style_name or "見出し 1" in style_name:
            md_lines.append(f"## {text}")
        elif "Heading 2" in style_name or "見出し 2" in style_name:
            md_lines.append(f"### {text}")
        elif "Heading 3" in style_name or "見出し 3" in style_name:
            md_lines.append(f"#### {text}")
        elif "Title" in style_name or "タイトル" in style_name:
            md_lines.append(f"## {text}")
        else:
            md_lines.append(text)
    
    return "\n".join(md_lines)


def convert_text_to_markdown(filepath: str) -> str:
    """テキストファイルをマークダウンに変換"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    md_lines = [f"# {Path(filepath).stem}\n"]
    md_lines.append(f"*変換元: {Path(filepath).name}*\n")
    md_lines.append(content)
    
    return "\n".join(md_lines)


def convert_file(filepath: str, output_dir: str) -> str:
    """ファイルをマークダウンに変換して保存"""
    path = Path(filepath)
    suffix = path.suffix.lower()
    
    converters = {
        '.pptx': convert_pptx_to_markdown,
        '.xlsx': convert_xlsx_to_markdown,
        '.xls': convert_xlsx_to_markdown,
        '.pdf': convert_pdf_to_markdown,
        '.docx': convert_docx_to_markdown,
        '.txt': convert_text_to_markdown,
        '.md': convert_text_to_markdown,
    }
    
    if suffix not in converters:
        print(f"未対応のファイル形式: {suffix}")
        return None
    
    print(f"変換中: {path.name}")
    md_content = converters[suffix](filepath)
    
    output_path = Path(output_dir) / f"{path.stem}.md"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    print(f"出力完了: {output_path}")
    return str(output_path)


def main():
    """メイン処理"""
    if len(sys.argv) < 3:
        print("使用方法: python convert_to_markdown.py <入力ファイルまたはディレクトリ> <出力ディレクトリ>")
        print("例: python convert_to_markdown.py /mnt/user-data/uploads /home/claude/pufu_work/01_markdown")
        print("")
        print("対応形式: .pptx, .xlsx, .xls, .pdf, .docx, .txt, .md")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_dir = sys.argv[2]
    
    input_path = Path(input_path)
    if input_path.is_file():
        files = [input_path]
    elif input_path.is_dir():
        extensions = ['.pptx', '.xlsx', '.xls', '.pdf', '.docx', '.txt', '.md']
        files = [f for f in input_path.iterdir() if f.suffix.lower() in extensions]
    else:
        print(f"ファイルまたはディレクトリが見つかりません: {input_path}")
        sys.exit(1)
    
    if not files:
        print("対象ファイルがありません")
        sys.exit(1)
    
    print(f"変換対象ファイル数: {len(files)}")
    
    converted = []
    for file in files:
        result = convert_file(str(file), output_dir)
        if result:
            converted.append(result)
    
    print(f"\n変換完了: {len(converted)}件")
    
    summary = {
        "convertedAt": datetime.now().isoformat(),
        "inputPath": str(input_path),
        "outputDir": output_dir,
        "files": converted
    }
    
    summary_path = Path(output_dir) / "_conversion_summary.json"
    with open(summary_path, 'w', encoding='utf-8') as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    
    print(f"サマリ出力: {summary_path}")


if __name__ == "__main__":
    main()
